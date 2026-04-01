# tools/format_tools.py
# =============================================================
# 📘 教学笔记：格式检查和调整工具
# =============================================================
# Agent 翻译完后可以自主检查输出效果，发现格式问题自己修复。
# inspect_output: 渲染输出文件的某一页为图片，让 Agent 看效果
# adjust_format: 调整指定段落的字号、加粗、对齐等格式
# =============================================================

import json
import os
import io
import base64
from typing import Dict, Optional

from core.agent_loop import BaseTool
from core.logger import get_logger

logger = get_logger("format_tools")


class InspectOutputTool(BaseTool):
    """
    📘 检查输出文件的视觉效果

    渲染输出文件的指定页为图片，返回 base64 让 Agent 看。
    Agent 对比原文和译文的视觉效果，自主决定是否需要调整。
    目前支持 PPT（通过 python-pptx + Pillow 简易渲染）和 PDF（PyMuPDF）。
    """

    name = "inspect_output"
    description = (
        "渲染输出文件的指定页为图片，让你检查翻译后的视觉效果。"
        "对比原文排版，检查字号是否合适、文字是否溢出、布局是否美观。"
        "如果发现问题，可以用 adjust_format 工具修复。"
    )
    parameters = {
        "type": "object",
        "properties": {
            "output_path": {
                "type": "string",
                "description": "输出文件路径",
            },
            "page_index": {
                "type": "integer",
                "description": "要检查的页码（0-based）",
            },
        },
        "required": ["output_path", "page_index"],
    }

    def execute(self, params: dict) -> str:
        output_path = params["output_path"]
        page_idx = params["page_index"]

        if not os.path.exists(output_path):
            return json.dumps({"error": f"文件不存在: {output_path}"}, ensure_ascii=False)

        ext = os.path.splitext(output_path)[1].lower()

        try:
            if ext == ".pptx":
                return self._inspect_pptx(output_path, page_idx)
            elif ext == ".pdf":
                return self._inspect_pdf(output_path, page_idx)
            elif ext == ".docx":
                return self._inspect_docx(output_path, page_idx)
            else:
                return json.dumps(
                    {"error": f"不支持检查此文件类型: {ext}"},
                    ensure_ascii=False,
                )
        except Exception as e:
            logger.error(f"检查输出失败: {e}")
            return json.dumps({"error": str(e)}, ensure_ascii=False)

    def _inspect_pptx(self, path: str, page_idx: int) -> str:
        """用 python-pptx 提取页面文本布局信息（无需 COM/LibreOffice）"""
        from pptx import Presentation
        from pptx.util import Emu

        prs = Presentation(path)
        if page_idx >= len(prs.slides):
            return json.dumps({"error": f"页码越界: {page_idx} >= {len(prs.slides)}"}, ensure_ascii=False)

        slide = prs.slides[page_idx]
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        shapes_info = []
        for shape in slide.shapes:
            info = {
                "name": shape.name,
                "left_pct": round(shape.left / slide_w * 100, 1) if slide_w else 0,
                "top_pct": round(shape.top / slide_h * 100, 1) if slide_h else 0,
                "width_pct": round(shape.width / slide_w * 100, 1) if slide_w else 0,
                "height_pct": round(shape.height / slide_h * 100, 1) if slide_h else 0,
            }
            if shape.has_text_frame:
                texts = []
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        font_sizes = []
                        for run in para.runs:
                            if run.font.size:
                                font_sizes.append(round(run.font.size.pt, 1))
                        texts.append({
                            "text": para_text[:80],
                            "font_sizes": font_sizes,
                            "char_count": len(para_text),
                        })
                if texts:
                    info["texts"] = texts
                    # 检测潜在问题
                    for t in texts:
                        if t["font_sizes"] and min(t["font_sizes"]) < 8:
                            info["warning"] = "字号过小（<8pt），可能不易阅读"
                        if t["char_count"] > 100 and info["width_pct"] < 30:
                            info["warning"] = "长文本在窄区域，可能溢出"
            if shape.has_table:
                info["type"] = "table"
                info["rows"] = len(shape.table.rows)
                info["cols"] = len(shape.table.columns)

            shapes_info.append(info)

        return json.dumps({
            "page_index": page_idx,
            "slide_size": f"{round(slide_w/914400, 1)}x{round(slide_h/914400, 1)} inches",
            "shapes_count": len(shapes_info),
            "shapes": shapes_info,
        }, ensure_ascii=False)

    def _inspect_pdf(self, path: str, page_idx: int) -> str:
        """用 PyMuPDF 渲染 PDF 页面为图片"""
        import fitz
        doc = fitz.open(path)
        if page_idx >= len(doc):
            doc.close()
            return json.dumps({"error": f"页码越界: {page_idx} >= {len(doc)}"}, ensure_ascii=False)

        page = doc[page_idx]
        pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
        img_bytes = pix.tobytes("jpeg", jpg_quality=80)
        doc.close()

        b64 = base64.b64encode(img_bytes).decode("utf-8")
        return json.dumps({
            "page_index": page_idx,
            "image_base64": b64,
            "width": pix.width,
            "height": pix.height,
        }, ensure_ascii=False)

    def _inspect_docx(self, path: str, page_idx: int) -> str:
        """检查 Word 文档的段落格式信息"""
        from docx import Document
        from docx.shared import Pt as DocxPt

        doc = Document(path)
        paragraphs = doc.paragraphs

        # Word 没有"页"的概念，用段落范围模拟
        # page_idx=0 → 前 20 段，page_idx=1 → 20-40 段，以此类推
        chunk_size = 20
        start = page_idx * chunk_size
        end = min(start + chunk_size, len(paragraphs))

        if start >= len(paragraphs):
            return json.dumps({
                "error": f"段落范围越界: 起始 {start} >= 总段落数 {len(paragraphs)}"
            }, ensure_ascii=False)

        para_infos = []
        for i in range(start, end):
            para = paragraphs[i]
            text = para.text.strip()
            if not text:
                continue

            font_sizes = []
            font_names = []
            for run in para.runs:
                if run.font.size:
                    font_sizes.append(round(run.font.size.pt, 1))
                if run.font.name:
                    font_names.append(run.font.name)

            info = {
                "index": i,
                "text": text[:80],
                "char_count": len(text),
                "style": para.style.name if para.style else None,
                "alignment": str(para.alignment) if para.alignment else None,
                "font_sizes": list(set(font_sizes)) if font_sizes else [],
                "font_names": list(set(font_names)) if font_names else [],
            }

            # 检测潜在问题
            if font_sizes and min(font_sizes) < 8:
                info["warning"] = "字号过小"
            if len(text) > 200:
                info["warning"] = "段落很长，检查是否需要分段"

            para_infos.append(info)

        # 也检查表格
        table_infos = []
        for t_idx, table in enumerate(doc.tables):
            table_infos.append({
                "table_index": t_idx,
                "rows": len(table.rows),
                "cols": len(table.columns),
            })

        return json.dumps({
            "page_index": page_idx,
            "paragraph_range": f"{start}-{end}",
            "total_paragraphs": len(paragraphs),
            "paragraphs": para_infos,
            "tables_count": len(table_infos),
            "tables": table_infos[:5],
        }, ensure_ascii=False)


class AdjustFormatTool(BaseTool):
    """
    📘 调整输出文件的格式

    Agent 检查后发现格式问题，用这个工具修复。
    支持调整字号、加粗、对齐、字体等。
    """

    name = "adjust_format"
    description = (
        "调整输出 PPT/Word 文件中指定段落的格式。"
        "可以调整字号、加粗、对齐方式、字体名称。"
        "用于修复翻译后的排版问题（如字号过小、文字溢出等）。"
    )
    parameters = {
        "type": "object",
        "properties": {
            "output_path": {
                "type": "string",
                "description": "输出文件路径",
            },
            "adjustments": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "page_index": {"type": "integer", "description": "页码（0-based）"},
                        "shape_name": {"type": "string", "description": "PPT 形状名称（从 inspect_output 获取）"},
                        "paragraph_index": {"type": "integer", "description": "Word 段落索引（从 inspect_output 获取）"},
                        "font_size_pt": {"type": "number", "description": "新字号（pt）"},
                        "bold": {"type": "boolean", "description": "是否加粗"},
                        "font_name": {"type": "string", "description": "字体名称"},
                    },
                    "required": ["page_index"],
                },
                "description": "格式调整列表",
            },
        },
        "required": ["output_path", "adjustments"],
    }

    def execute(self, params: dict) -> str:
        output_path = params["output_path"]
        adjustments = params["adjustments"]

        if not os.path.exists(output_path):
            return json.dumps({"error": f"文件不存在: {output_path}"}, ensure_ascii=False)

        ext = os.path.splitext(output_path)[1].lower()

        try:
            if ext == ".pptx":
                return self._adjust_pptx(output_path, adjustments)
            elif ext == ".docx":
                return self._adjust_docx(output_path, adjustments)
            else:
                return json.dumps(
                    {"error": f"暂不支持调整此文件类型: {ext}"},
                    ensure_ascii=False,
                )
        except Exception as e:
            logger.error(f"格式调整失败: {e}")
            return json.dumps({"error": str(e)}, ensure_ascii=False)

    def _adjust_pptx(self, path: str, adjustments: list) -> str:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(path)
        adjusted_count = 0

        for adj in adjustments:
            page_idx = adj["page_index"]
            if page_idx >= len(prs.slides):
                continue

            slide = prs.slides[page_idx]
            shape_name = adj.get("shape_name")
            font_size = adj.get("font_size_pt")
            bold = adj.get("bold")
            font_name = adj.get("font_name")

            # 找到目标 shape
            target_shapes = []
            if shape_name:
                for shape in slide.shapes:
                    if shape.name == shape_name:
                        target_shapes.append(shape)
            else:
                # 没指定 shape_name，调整该页所有文本形状
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        target_shapes.append(shape)

            for shape in target_shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if font_size is not None:
                            run.font.size = Pt(font_size)
                        if bold is not None:
                            run.font.bold = bold
                        if font_name is not None:
                            run.font.name = font_name
                        adjusted_count += 1

        prs.save(path)
        return json.dumps({
            "success": True,
            "adjusted_runs": adjusted_count,
        }, ensure_ascii=False)

    def _adjust_docx(self, path: str, adjustments: list) -> str:
        """调整 Word 文档的段落格式"""
        from docx import Document
        from docx.shared import Pt as DocxPt
        from docx.oxml.ns import qn

        doc = Document(path)
        paragraphs = doc.paragraphs
        adjusted_count = 0

        for adj in adjustments:
            para_index = adj.get("paragraph_index")
            font_size = adj.get("font_size_pt")
            bold = adj.get("bold")
            font_name = adj.get("font_name")

            if para_index is not None and para_index < len(paragraphs):
                target_paras = [paragraphs[para_index]]
            else:
                # page_index 模拟：每 20 段一页。page_index=-1 表示全部段落
                page_idx = adj.get("page_index", -1)
                if page_idx == -1:
                    target_paras = list(paragraphs)
                else:
                    chunk_size = 20
                    start = page_idx * chunk_size
                    end = min(start + chunk_size, len(paragraphs))
                    target_paras = paragraphs[start:end]

            for para in target_paras:
                for run in para.runs:
                    if not run.text:
                        continue  # 跳过空 Run
                    if font_size is not None:
                        run.font.size = DocxPt(font_size)
                    if bold is not None:
                        run.font.bold = bold
                    if font_name is not None:
                        run.font.name = font_name
                        # 📘 同时设置东亚字体，确保中文也生效
                        rPr = run._element.get_or_add_rPr()
                        rFonts = rPr.find(qn('w:rFonts'))
                        if rFonts is None:
                            from lxml import etree
                            rFonts = etree.SubElement(rPr, qn('w:rFonts'))
                        rFonts.set(qn('w:ascii'), font_name)
                        rFonts.set(qn('w:hAnsi'), font_name)
                        rFonts.set(qn('w:eastAsia'), font_name)
                    adjusted_count += 1

            # 📘 也处理表格中的文字
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.paragraphs:
                            for run in para.runs:
                                if not run.text:
                                    continue
                                if font_size is not None:
                                    run.font.size = DocxPt(font_size)
                                if bold is not None:
                                    run.font.bold = bold
                                if font_name is not None:
                                    run.font.name = font_name
                                    rPr = run._element.get_or_add_rPr()
                                    rFonts = rPr.find(qn('w:rFonts'))
                                    if rFonts is None:
                                        from lxml import etree
                                        rFonts = etree.SubElement(rPr, qn('w:rFonts'))
                                    rFonts.set(qn('w:ascii'), font_name)
                                    rFonts.set(qn('w:hAnsi'), font_name)
                                    rFonts.set(qn('w:eastAsia'), font_name)
                                adjusted_count += 1

        doc.save(path)
        return json.dumps({
            "success": True,
            "adjusted_runs": adjusted_count,
        }, ensure_ascii=False)
