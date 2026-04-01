# tools/layout_tools_v2.py
# =============================================================
# 📘 教学笔记：高级排版工具（v2 — Agent 自主格式优化）
# =============================================================
# 这些工具让 Agent 能真正"看到"输出效果并精确调整排版。
#
# 核心思路：
# 1. render_slide: 用 COM 把 PPT 幻灯片渲染为真实图片（Agent 能看到）
# 2. enable_autofit: 给文本框设置 PPT 原生的"缩小文字以适应"
# 3. compare_layout: 对比原文和译文的视觉布局差异
# 4. smart_resize: 智能调整字号（基于文本框尺寸和文本长度计算）
# =============================================================

import json
import os
import io
import base64
import math
from typing import Dict, List, Optional

from core.agent_loop import BaseTool
from core.logger import get_logger

logger = get_logger("layout_tools_v2")


class RenderSlideTool(BaseTool):
    """
    📘 渲染 PPT 幻灯片为真实图片

    用 Windows COM（PowerPoint）把幻灯片导出为 PNG。
    Agent 能看到真实的渲染效果，包括字体渲染、文字溢出等。
    比 inspect_output 的纯数据分析更直观。

    需要 Windows + Office 安装。
    """

    name = "render_slide"
    description = (
        "把 PPT 的指定幻灯片渲染为真实图片（PNG），让你看到实际效果。"
        "包括字体渲染、文字是否溢出、整体布局是否美观。"
        "需要 Windows 系统且安装了 PowerPoint。"
    )
    parameters = {
        "type": "object",
        "properties": {
            "pptx_path": {
                "type": "string",
                "description": "PPT 文件路径",
            },
            "slide_index": {
                "type": "integer",
                "description": "幻灯片索引（0-based）",
            },
        },
        "required": ["pptx_path", "slide_index"],
    }

    def execute(self, params: dict) -> str:
        pptx_path = params["pptx_path"]
        slide_idx = params["slide_index"]

        if not os.path.exists(pptx_path):
            return json.dumps({"error": f"文件不存在: {pptx_path}"}, ensure_ascii=False)

        try:
            img_b64, width, height = self._render_com(pptx_path, slide_idx)
            return json.dumps({
                "slide_index": slide_idx,
                "image_base64": img_b64,
                "width": width,
                "height": height,
                "method": "COM",
            }, ensure_ascii=False)
        except Exception as e:
            logger.warning(f"COM 渲染失败: {e}，降级为布局分析")
            # 降级：返回文本布局信息（不需要 COM）
            return self._fallback_inspect(pptx_path, slide_idx)

    def _render_com(self, pptx_path: str, slide_idx: int) -> tuple:
        """用 COM 渲染幻灯片为 PNG"""
        import win32com.client
        import tempfile
        from PIL import Image

        abs_path = os.path.abspath(pptx_path)
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = False

        try:
            prs = ppt_app.Presentations.Open(abs_path, WithWindow=False)
            slide = prs.Slides(slide_idx + 1)  # COM 是 1-based

            # 导出为临时 PNG
            tmp_dir = tempfile.mkdtemp()
            tmp_path = os.path.join(tmp_dir, f"slide_{slide_idx}.png")
            slide.Export(tmp_path, "PNG", 1920, 1080)
            prs.Close()

            # 读取并转 base64
            img = Image.open(tmp_path)
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=85)
            b64 = base64.b64encode(buf.getvalue()).decode("utf-8")

            w, h = img.size
            # 清理临时文件
            os.remove(tmp_path)
            os.rmdir(tmp_dir)

            return b64, w, h
        finally:
            try:
                ppt_app.Quit()
            except Exception:
                pass

    def _fallback_inspect(self, pptx_path: str, slide_idx: int) -> str:
        """降级：返回文本布局信息"""
        from pptx import Presentation
        prs = Presentation(pptx_path)
        if slide_idx >= len(prs.slides):
            return json.dumps({"error": "页码越界"}, ensure_ascii=False)

        slide = prs.slides[slide_idx]
        info = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        sizes = [round(r.font.size.pt, 1) for r in para.runs if r.font.size]
                        info.append({"text": text[:60], "font_sizes": sizes})
        return json.dumps({
            "slide_index": slide_idx,
            "method": "fallback_text_analysis",
            "texts": info,
        }, ensure_ascii=False)


class EnableAutofitTool(BaseTool):
    """
    📘 给 PPT 文本框启用原生"缩小文字以适应"

    PowerPoint 有原生的 normAutofit 功能：
    当文字超出文本框时，自动缩小字号直到文字完全适配。
    这比我们手动计算字号更准确（PowerPoint 知道精确的字体度量）。

    通过直接操作 XML 设置 <a:normAutofit/>，
    PowerPoint 打开文件时自动计算最佳字号。
    """

    name = "enable_autofit"
    description = (
        "给 PPT 指定页面的所有文本框启用 PowerPoint 原生的"
        "'缩小文字以适应'功能。PowerPoint 打开时会自动计算最佳字号，"
        "确保文字不溢出。这是最可靠的防溢出方案。"
    )
    parameters = {
        "type": "object",
        "properties": {
            "pptx_path": {
                "type": "string",
                "description": "PPT 文件路径",
            },
            "slide_indices": {
                "type": "array",
                "items": {"type": "integer"},
                "description": "要启用 autofit 的幻灯片索引列表（0-based）。空数组=全部页面。",
            },
        },
        "required": ["pptx_path"],
    }

    def execute(self, params: dict) -> str:
        pptx_path = params["pptx_path"]
        slide_indices = params.get("slide_indices", [])

        if not os.path.exists(pptx_path):
            return json.dumps({"error": f"文件不存在: {pptx_path}"}, ensure_ascii=False)

        try:
            from pptx import Presentation
            from lxml import etree

            prs = Presentation(pptx_path)
            nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            count = 0

            target_slides = []
            if slide_indices:
                for idx in slide_indices:
                    if idx < len(prs.slides):
                        target_slides.append(prs.slides[idx])
            else:
                target_slides = list(prs.slides)

            for slide in target_slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    # 📘 直接操作 XML：设置 <a:bodyPr> 中的 <a:normAutofit/>
                    txBody = shape.text_frame._txBody
                    bodyPr = txBody.find(".//a:bodyPr", nsmap)
                    if bodyPr is None:
                        continue

                    # 移除已有的 autofit 设置
                    for child in bodyPr:
                        tag = etree.QName(child.tag).localname
                        if tag in ("noAutofit", "spAutoFit", "normAutofit"):
                            bodyPr.remove(child)

                    # 添加 normAutofit（缩小文字以适应）
                    norm = etree.SubElement(
                        bodyPr,
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}normAutofit"
                    )
                    count += 1

            prs.save(pptx_path)
            return json.dumps({
                "success": True,
                "shapes_updated": count,
                "slides_processed": len(target_slides),
            }, ensure_ascii=False)

        except Exception as e:
            logger.error(f"启用 autofit 失败: {e}")
            return json.dumps({"error": str(e)}, ensure_ascii=False)


class CompareLayoutTool(BaseTool):
    """
    📘 对比原文和译文的视觉布局

    把原文和译文的同一页并排对比，分析布局差异。
    Agent 用这个工具判断译文排版是否需要调整。
    """

    name = "compare_layout"
    description = (
        "对比原文和译文同一页的布局差异。"
        "分析每个文本框的位置、尺寸、字号变化，"
        "找出可能的排版问题（文字溢出、字号过小、布局偏移）。"
    )
    parameters = {
        "type": "object",
        "properties": {
            "original_path": {
                "type": "string",
                "description": "原文文件路径",
            },
            "translated_path": {
                "type": "string",
                "description": "译文文件路径",
            },
            "slide_index": {
                "type": "integer",
                "description": "要对比的幻灯片索引（0-based）",
            },
        },
        "required": ["original_path", "translated_path", "slide_index"],
    }

    def execute(self, params: dict) -> str:
        orig_path = params["original_path"]
        trans_path = params["translated_path"]
        slide_idx = params["slide_index"]

        try:
            from pptx import Presentation

            orig_prs = Presentation(orig_path)
            trans_prs = Presentation(trans_path)

            if slide_idx >= len(orig_prs.slides) or slide_idx >= len(trans_prs.slides):
                return json.dumps({"error": "页码越界"}, ensure_ascii=False)

            orig_slide = orig_prs.slides[slide_idx]
            trans_slide = trans_prs.slides[slide_idx]

            diffs = []
            orig_shapes = {s.name: s for s in orig_slide.shapes}
            trans_shapes = {s.name: s for s in trans_slide.shapes}

            for name, orig_shape in orig_shapes.items():
                if name not in trans_shapes:
                    continue
                trans_shape = trans_shapes[name]

                if not orig_shape.has_text_frame or not trans_shape.has_text_frame:
                    continue

                orig_text = orig_shape.text_frame.text.strip()
                trans_text = trans_shape.text_frame.text.strip()

                if not orig_text and not trans_text:
                    continue

                # 字号对比
                orig_sizes = []
                trans_sizes = []
                for para in orig_shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            orig_sizes.append(run.font.size.pt)
                for para in trans_shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            trans_sizes.append(run.font.size.pt)

                diff = {
                    "shape_name": name,
                    "original_text": orig_text[:50],
                    "translated_text": trans_text[:50],
                    "original_char_count": len(orig_text),
                    "translated_char_count": len(trans_text),
                    "length_ratio": round(len(trans_text) / max(len(orig_text), 1), 2),
                }

                if orig_sizes and trans_sizes:
                    diff["original_font_size"] = round(min(orig_sizes), 1)
                    diff["translated_font_size"] = round(min(trans_sizes), 1)
                    size_change = round(min(trans_sizes) - min(orig_sizes), 1)
                    if size_change < -2:
                        diff["issue"] = f"字号缩小了 {abs(size_change)}pt"
                    if min(trans_sizes) < 8:
                        diff["issue"] = "字号过小，可能不易阅读"

                if len(trans_text) > len(orig_text) * 2:
                    diff["issue"] = "译文长度是原文的2倍以上，可能溢出"

                diffs.append(diff)

            issues = [d for d in diffs if "issue" in d]
            return json.dumps({
                "slide_index": slide_idx,
                "total_shapes": len(diffs),
                "issues_found": len(issues),
                "details": diffs,
            }, ensure_ascii=False)

        except Exception as e:
            logger.error(f"布局对比失败: {e}")
            return json.dumps({"error": str(e)}, ensure_ascii=False)


class SmartResizeTool(BaseTool):
    """
    📘 智能字号调整

    根据文本框尺寸和文本长度，计算最佳字号。
    不是简单的按比例缩放，而是考虑：
    - 文本框的宽度和高度
    - 英文字符的平均宽度（比中文窄）
    - 换行后的行数
    - 最小可读字号下限
    """

    name = "smart_resize"
    description = (
        "智能计算并设置最佳字号。根据文本框尺寸和文本长度，"
        "计算能让文字完全适配且美观的字号。"
        "比手动设置更精确，比 enable_autofit 更可控。"
    )
    parameters = {
        "type": "object",
        "properties": {
            "pptx_path": {
                "type": "string",
                "description": "PPT 文件路径",
            },
            "slide_index": {
                "type": "integer",
                "description": "幻灯片索引（0-based）",
            },
            "shape_name": {
                "type": "string",
                "description": "形状名称（可选，不指定则处理该页所有文本框）",
            },
            "min_font_size": {
                "type": "number",
                "description": "最小字号下限（pt），默认 8",
            },
        },
        "required": ["pptx_path", "slide_index"],
    }

    def execute(self, params: dict) -> str:
        pptx_path = params["pptx_path"]
        slide_idx = params["slide_index"]
        shape_name = params.get("shape_name")
        min_size = params.get("min_font_size", 8)

        try:
            from pptx import Presentation
            from pptx.util import Pt, Emu

            prs = Presentation(pptx_path)
            if slide_idx >= len(prs.slides):
                return json.dumps({"error": "页码越界"}, ensure_ascii=False)

            slide = prs.slides[slide_idx]
            results = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape_name and shape.name != shape_name:
                    continue

                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # 文本框尺寸（EMU -> pt: 1pt = 12700 EMU）
                box_width_pt = shape.width / 12700
                box_height_pt = shape.height / 12700

                # 获取当前字号
                current_sizes = []
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            current_sizes.append(run.font.size.pt)

                if not current_sizes:
                    continue

                current_size = min(current_sizes)

                # 📘 计算最佳字号
                # 英文字符平均宽度约 0.6 * font_size
                # 中文字符平均宽度约 1.0 * font_size
                avg_char_width_ratio = 0.6  # 英文为主
                for ch in text:
                    if '\u4e00' <= ch <= '\u9fff':
                        avg_char_width_ratio = 0.8  # 中英混合
                        break

                # 每行能放多少字符
                chars_per_line = box_width_pt / (current_size * avg_char_width_ratio)
                # 需要多少行
                lines_needed = math.ceil(len(text) / max(chars_per_line, 1))
                # 行高约 1.2 * font_size
                height_needed = lines_needed * current_size * 1.2

                if height_needed <= box_height_pt:
                    # 当前字号够用，不需要调整
                    results.append({
                        "shape": shape.name,
                        "action": "no_change",
                        "current_size": current_size,
                    })
                    continue

                # 需要缩小：二分查找最佳字号
                best_size = current_size
                lo, hi = min_size, current_size
                for _ in range(20):  # 最多迭代 20 次
                    mid = (lo + hi) / 2
                    cpl = box_width_pt / (mid * avg_char_width_ratio)
                    ln = math.ceil(len(text) / max(cpl, 1))
                    h = ln * mid * 1.2
                    if h <= box_height_pt:
                        best_size = mid
                        lo = mid
                    else:
                        hi = mid
                    if hi - lo < 0.5:
                        break

                best_size = max(round(best_size, 1), min_size)

                # 应用新字号
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(best_size)

                results.append({
                    "shape": shape.name,
                    "action": "resized",
                    "original_size": current_size,
                    "new_size": best_size,
                    "text_preview": text[:40],
                })

            prs.save(pptx_path)
            resized = [r for r in results if r["action"] == "resized"]
            return json.dumps({
                "success": True,
                "total_shapes": len(results),
                "resized": len(resized),
                "details": results,
            }, ensure_ascii=False)

        except Exception as e:
            logger.error(f"智能调整失败: {e}")
            return json.dumps({"error": str(e)}, ensure_ascii=False)
