# translator/pptx_writer.py
import re
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Dict, Any, Optional
from translator.format_engine import FormatEngine
from core.logger import get_logger

# =============================================================
# 📘 教学笔记：PPT 文档生成器（克隆策略）
# =============================================================
# 和 docx_writer 一样，我们采用"克隆原文件 → 原地替换"的策略：
#   1. 用 python-pptx 打开原始 .pptx 文件
#   2. 按 key 定位到对应的段落
#   3. 替换 Run 文本，映射字体
#   4. 保存为新文件
#
# 这样所有原始格式（动画、母版、配色、图片位置等）都完整保留。
#
# 和 Word 的区别：
#   - PPT 的形状是二维布局（不是线性文档流）
#   - 需要按 slide → shape → paragraph 三级定位
#   - 组合形状需要递归进入
#   - 表格在 shape.table 里，不是 doc.tables
# =============================================================

logger = get_logger("pptx_writer")

RUN_TAG_PATTERN = re.compile(r'<r(\d+)>(.*?)</r\1>', re.DOTALL)


def _parse_tagged_text(text: str) -> Optional[Dict[int, str]]:
    """解析带标记的译文，返回 {Run编号: 译文} 字典。"""
    matches = RUN_TAG_PATTERN.findall(text)
    if not matches:
        return None
    return {int(idx): content for idx, content in matches}


def _remap_run_font(run, format_engine: FormatEngine):
    """替换 Run 的字体名，其他格式保留。"""
    original_font = run.font.name
    if original_font:
        new_font = format_engine.resolve_font(original_font)
        if new_font:
            run.font.name = new_font


def _shrink_runs_if_needed(paragraph, original_text: str, translated_text: str):
    """
    📘 教学笔记：按长度比例缩小字号
    auto_size 依赖 PowerPoint 打开时才计算，不可靠。
    我们直接在写入时计算：如果译文比原文长，就按比例缩小所有 Run 的字号。

    中文字符按 2 个宽度单位计算（一个汉字约等于两个英文字母宽度），
    这样中→英翻译时的比例更准确。

    设置最小字号下限（6pt），避免缩得太小看不清。
    """
    def _visual_len(text: str) -> float:
        """估算文本的视觉宽度（中文字符算2，其他算1）"""
        width = 0
        for ch in text:
            if '\u4e00' <= ch <= '\u9fff' or '\u3000' <= ch <= '\u303f' or '\uff00' <= ch <= '\uffef':
                width += 2
            else:
                width += 1
        return width

    # 清理标记后计算
    clean_original = RUN_TAG_PATTERN.sub(r'\2', original_text)
    clean_translated = RUN_TAG_PATTERN.sub(r'\2', translated_text)

    orig_len = _visual_len(clean_original)
    trans_len = _visual_len(clean_translated)

    if orig_len == 0 or trans_len <= orig_len:
        return  # 没变长，不需要缩

    ratio = orig_len / trans_len  # < 1.0，表示需要缩小

    MIN_FONT_SIZE = Pt(6)

    for run in paragraph.runs:
        if run.font.size and run.font.size > MIN_FONT_SIZE:
            new_size = int(run.font.size * ratio)
            run.font.size = max(new_size, MIN_FONT_SIZE)


def _replace_paragraph_text(paragraph, translated_text: str, is_tagged: bool,
                            format_engine: FormatEngine, original_text: str = ""):
    """
    替换一个段落的文本（通用逻辑）。
    和 docx_writer 的逻辑基本一致。
    """
    original_runs = paragraph.runs

    if len(original_runs) == 0:
        clean = RUN_TAG_PATTERN.sub(r'\2', translated_text)
        paragraph.text = clean
        return

    # ---- 多 Run + 带标记：逐 Run 替换 ----
    if is_tagged and len(original_runs) > 1:
        run_texts = _parse_tagged_text(translated_text)
        if run_texts is not None:
            non_empty_indices = [ri for ri, run in enumerate(original_runs) if run.text]
            success = True
            for tag_idx, text in run_texts.items():
                if tag_idx < len(non_empty_indices):
                    actual_idx = non_empty_indices[tag_idx]
                    original_runs[actual_idx].text = text
                    _remap_run_font(original_runs[actual_idx], format_engine)
                else:
                    logger.warning(f"标记 r{tag_idx} 超出 Run 范围，降级处理")
                    success = False
                    break

            if success:
                tagged_indices = set()
                for tag_idx in run_texts.keys():
                    if tag_idx < len(non_empty_indices):
                        tagged_indices.add(non_empty_indices[tag_idx])
                for ri, run in enumerate(original_runs):
                    if run.text and ri not in tagged_indices:
                        run.text = ""
                # 标记替换成功后，也要缩小字号
                _shrink_runs_if_needed(paragraph, original_text, translated_text)
                return

        logger.warning("标记解析失败，降级为整段替换")

    # ---- 单 Run 或降级：整段替换 ----
    clean_text = RUN_TAG_PATTERN.sub(r'\2', translated_text)
    original_runs[0].text = clean_text
    _remap_run_font(original_runs[0], format_engine)
    for run in original_runs[1:]:
        run.text = ""

    # 替换完成后，按长度比例缩小字号
    _shrink_runs_if_needed(paragraph, original_text, translated_text)


def _resolve_paragraph_and_shape(shapes, key: str):
    """
    根据 key 在形状集合中定位到具体的段落对象，同时返回所属的 shape。

    返回 (paragraph, shape) 或 (None, None)。
    """
    parts = key.split("_")
    idx = 1  # 跳过 "s0"

    current_shapes = shapes
    while idx < len(parts) and parts[idx].startswith("g"):
        group_idx = int(parts[idx][1:])
        shape_list = list(current_shapes)
        if group_idx >= len(shape_list):
            return None, None
        group_shape = shape_list[group_idx]
        if group_shape.shape_type != MSO_SHAPE_TYPE.GROUP:
            return None, None
        current_shapes = group_shape.shapes
        idx += 1

    if idx >= len(parts) or not parts[idx].startswith("sh"):
        return None, None
    shape_idx = int(parts[idx][2:])
    shape_list = list(current_shapes)
    if shape_idx >= len(shape_list):
        return None, None
    shape = shape_list[shape_idx]
    idx += 1

    if idx >= len(parts):
        return None, None

    # 表格：t{row}_{col}[_p{para}]
    if parts[idx].startswith("t"):
        row_idx = int(parts[idx][1:])
        idx += 1
        if idx >= len(parts):
            return None, None
        col_idx = int(parts[idx])
        idx += 1
        para_idx = 0
        if idx < len(parts) and parts[idx].startswith("p"):
            para_idx = int(parts[idx][1:])

        if not shape.has_table:
            return None, None
        table = shape.table
        if row_idx >= len(table.rows) or col_idx >= len(table.rows[row_idx].cells):
            return None, None
        cell = table.rows[row_idx].cells[col_idx]
        paras = cell.text_frame.paragraphs
        if para_idx >= len(paras):
            return None, None
        return paras[para_idx], shape

    # 普通文本：p{para}
    if parts[idx].startswith("p"):
        para_idx = int(parts[idx][1:])
        if not shape.has_text_frame:
            return None, None
        paras = shape.text_frame.paragraphs
        if para_idx >= len(paras):
            return None, None
        return paras[para_idx], shape

    return None, None


def write_pptx(
    parsed_data: Dict[str, Any],
    translations: Dict[str, str],
    output_path: str,
    format_engine: FormatEngine,
    source_path: str = None,
):
    """
    基于原 PPT 生成翻译后的文件（克隆 + 原地替换）。

    📘 教学笔记：自动缩小字体防溢出
    翻译后文本往往比原文长（尤其中→英），容易撑破文本框导致换行。
    我们的策略：给所有被翻译过的文本框开启 PPT 原生的
    "Shrink text on overflow"（溢出时缩小字体）功能。
    PowerPoint 会自动计算最合适的字号，保证文字不溢出。
    表格单元格不需要处理（单元格会自动扩展高度）。
    """
    if not source_path:
        raise ValueError("source_path 不能为空，需要原 PPT 来克隆格式")

    logger.info(f"开始生成 PPT: {output_path}")
    prs = Presentation(source_path)

    replaced_count = 0

    for item in parsed_data["items"]:
        key = item["key"]
        if key not in translations:
            continue

        slide_idx = int(key.split("_")[0][1:])
        if slide_idx >= len(prs.slides):
            logger.warning(f"幻灯片索引越界 key={key}，跳过")
            continue

        slide = prs.slides[slide_idx]
        paragraph, shape = _resolve_paragraph_and_shape(slide.shapes, key)
        if paragraph is None:
            logger.warning(f"无法定位 key={key}，跳过")
            continue

        translated_text = translations[key]
        is_tagged = item.get("tagged_text", False)
        original_text = item.get("full_text", "")

        _replace_paragraph_text(paragraph, translated_text, is_tagged,
                                format_engine, original_text)
        replaced_count += 1

    prs.save(output_path)
    logger.info(f"PPT 生成完成: {output_path}（替换了 {replaced_count} 个翻译单元）")
