# translator/pptx_parser.py
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Dict, Any, Optional
from core.logger import get_logger

# =============================================================
# 📘 教学笔记：PPT 文档解析器
# =============================================================
# PPT 的结构和 Word 完全不同：
#   Word: 线性结构 → 段落1 → 段落2 → 表格 → 段落3
#   PPT:  树形结构 → 幻灯片 → 形状(Shape) → 文本框/表格/组合
#
# python-pptx 的对象层次：
#   Presentation
#     └─ Slides[]
#          └─ Shapes[]
#               ├─ TextFrame → Paragraphs[] → Runs[]
#               ├─ Table → Cells[] → TextFrame → ...
#               └─ GroupShapes[] → (递归包含更多 Shape)
#
# 我们的解析策略：
#   1. 遍历每张幻灯片的所有形状
#   2. 递归处理组合形状（GroupShape）
#   3. 对每个含文本的形状，按段落+Run 提取
#   4. 表格单独处理（遍历单元格）
#   5. 跳过图片、视频、音频等媒体形状
#
# Key 命名体系：
#   s{slide}_sh{shape}_p{para}           → 普通形状的段落
#   s{slide}_sh{shape}_t{row}_{col}      → 表格单元格（单段落）
#   s{slide}_sh{shape}_t{row}_{col}_p{para} → 表格单元格（多段落）
#   s{slide}_g{group}_sh{shape}_p{para}  → 组合形状内的段落
# =============================================================

logger = get_logger("pptx_parser")


def _extract_run_format(run) -> Dict[str, Any]:
    """提取一个 Run 的格式信息"""
    font = run.font
    # 📘 教学笔记：PPT 颜色类型比 Word 复杂
    # Word 基本都是 RGB，但 PPT 有多种颜色类型：
    #   - RGBColor: 直接的 RGB 值，可以访问 .rgb
    #   - SchemeColor: 主题配色方案中的颜色（如"强调色1"）
    #   - ThemeColor: 主题颜色
    # 只有 RGBColor 类型才能安全访问 .rgb，其他类型会抛异常。
    font_color = None
    try:
        if font.color and font.color.rgb:
            font_color = str(font.color.rgb)
    except (AttributeError, TypeError):
        pass

    fmt = {
        "bold": font.bold,
        "italic": font.italic,
        "underline": font.underline,
        "font_name": font.name,
        "font_size": font.size.pt if font.size else None,
        "font_color": font_color,
    }
    return fmt


def _build_tagged_text(runs_data: List[Dict]) -> str:
    """把多个 Run 的文本用标记包裹"""
    parts = []
    for i, run in enumerate(runs_data):
        parts.append(f"<r{i}>{run['text']}</r{i}>")
    return "".join(parts)


def _parse_paragraph(para, key: str) -> Optional[Dict[str, Any]]:
    """
    解析一个段落，返回翻译单元。
    和 docx_parser 的逻辑基本一致，复用 Run 标记策略。
    """
    raw_text = para.text.strip()
    if not raw_text:
        return None

    runs = []
    for run in para.runs:
        if not run.text:
            continue
        runs.append({
            "text": run.text,
            "format": _extract_run_format(run),
        })

    if not runs:
        return None

    # 判断是否需要 Run 标记（多 Run 且格式不同）
    needs_tagging = False
    if len(runs) > 1:
        first_fmt = runs[0]["format"]
        for r in runs[1:]:
            if r["format"] != first_fmt:
                needs_tagging = True
                break

    full_text = _build_tagged_text(runs) if needs_tagging else raw_text

    return {
        "key": key,
        "type": "slide_text",
        "full_text": full_text,
        "runs": runs,
        "tagged_text": needs_tagging,
    }


def _parse_table_shape(shape, slide_idx: int, shape_idx: int) -> List[Dict[str, Any]]:
    """解析表格形状，返回所有单元格的翻译单元"""
    items = []
    table = shape.table
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            paragraphs = cell.text_frame.paragraphs
            if len(paragraphs) == 1:
                key = f"s{slide_idx}_sh{shape_idx}_t{row_idx}_{col_idx}"
                result = _parse_paragraph(paragraphs[0], key)
                if result:
                    result["type"] = "table_cell"
                    items.append(result)
            else:
                for para_idx, para in enumerate(paragraphs):
                    key = f"s{slide_idx}_sh{shape_idx}_t{row_idx}_{col_idx}_p{para_idx}"
                    result = _parse_paragraph(para, key)
                    if result:
                        result["type"] = "table_cell"
                        items.append(result)
    return items


def _parse_text_shape(shape, key_prefix: str) -> List[Dict[str, Any]]:
    """解析含文本框的形状，返回段落翻译单元"""
    items = []
    if not shape.has_text_frame:
        return items
    for para_idx, para in enumerate(shape.text_frame.paragraphs):
        key = f"{key_prefix}_p{para_idx}"
        result = _parse_paragraph(para, key)
        if result:
            items.append(result)
    return items


def _is_media_shape(shape) -> bool:
    """判断是否为媒体形状（图片/视频/音频），这些不需要翻译"""
    # 📘 教学笔记：跳过媒体类型
    # MSO_SHAPE_TYPE 枚举定义了所有形状类型
    # 我们只跳过纯媒体，其他都尝试提取文本
    media_types = {
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.MEDIA,
    }
    try:
        return shape.shape_type in media_types
    except Exception:
        return False


def _parse_shapes(shapes, slide_idx: int, group_prefix: str = "") -> List[Dict[str, Any]]:
    """
    递归解析形状集合（支持 GroupShape 嵌套）。

    📘 教学笔记：递归处理组合形状
    PPT 里的"组合"操作会把多个形状包在一个 GroupShape 里。
    GroupShape 本身没有文本，但它的子形状有。
    我们用递归来"拆开"所有层级的组合。
    """
    items = []
    for shape_idx, shape in enumerate(shapes):
        if _is_media_shape(shape):
            continue

        prefix = f"s{slide_idx}_{group_prefix}sh{shape_idx}"

        # 组合形状：递归处理
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            nested_prefix = f"{group_prefix}g{shape_idx}_"
            items.extend(_parse_shapes(shape.shapes, slide_idx, nested_prefix))
            continue

        # 表格形状
        if shape.has_table:
            items.extend(_parse_table_shape(shape, slide_idx, shape_idx))
            continue

        # 含文本的形状（文本框、占位符、自选图形等）
        if shape.has_text_frame:
            items.extend(_parse_text_shape(shape, prefix))

    return items


def parse_pptx(filepath: str) -> Dict[str, Any]:
    """
    解析 PPT 文档，返回所有翻译单元。

    返回格式与 docx_parser.parse_docx() 一致：
    {
        "items": [
            {"key": "s0_sh1_p0", "type": "slide_text", ...},
            {"key": "s0_sh2_t0_1", "type": "table_cell", ...},
            ...
        ]
    }
    """
    logger.info(f"开始解析 PPT: {filepath}")
    prs = Presentation(filepath)

    items = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_items = _parse_shapes(slide.shapes, slide_idx)
        items.extend(slide_items)

    # 统计
    text_count = sum(1 for i in items if i["type"] == "slide_text")
    cell_count = sum(1 for i in items if i["type"] == "table_cell")
    tagged_count = sum(1 for i in items if i.get("tagged_text"))

    logger.info(f"PPT 解析完成: {text_count} 个文本段落，{cell_count} 个表格单元格，"
                f"{tagged_count} 个含格式标记，共 {len(prs.slides)} 张幻灯片")

    return {"items": items}
